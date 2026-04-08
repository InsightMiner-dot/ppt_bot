import os
import hashlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
from langchain_chroma import Chroma


# =========================
# ✅ CONFIG
# =========================
PERSIST_DIR = "chroma_db"

AZURE_OPENAI_API_KEY = "YOUR_KEY"
AZURE_OPENAI_ENDPOINT = "YOUR_ENDPOINT"

EMBEDDING_DEPLOYMENT = "text-embedding-3-small"
CHAT_DEPLOYMENT = "gpt-4o-mini"
API_VERSION = "2024-02-01"


# =========================
# ✅ VERSION
# =========================
def get_file_hash(file_path):
    with open(file_path, "rb") as f:
        return hashlib.md5(f.read()).hexdigest()


# =========================
# ✅ NORMALIZATION (CRITICAL)
# =========================
def normalize(text):
    return (
        text.lower()
        .replace("&", "and")
        .replace("-", "")
        .replace("  ", " ")
        .strip()
    )


# =========================
# ✅ EXTRACTION
# =========================
def extract_ppt(file_path):
    prs = Presentation(file_path)
    file_name = os.path.basename(file_path)
    version = get_file_hash(file_path)

    docs = []

    for i, slide in enumerate(prs.slides):
        comments = []
        slide_title = None
        chart_count = 0

        for shape in slide.shapes:

            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    comments.append(text)
                    if not slide_title:
                        slide_title = text

            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_count += 1

        docs.append({
            "content": "\n".join(comments),
            "metadata": {
                "file_name": file_name,
                "version": version,
                "slide_number": i + 1,
                "slide_title": slide_title,
                "chart_count": chart_count,
                "source": f"{file_name}_slide_{i+1}"
            }
        })

    return docs


# =========================
# ✅ TO DOCUMENT
# =========================
def to_docs(docs):
    return [Document(page_content=d["content"], metadata=d["metadata"]) for d in docs]


# =========================
# ✅ CHUNK ID
# =========================
def chunk_id(file, slide, idx, version):
    return hashlib.md5(f"{file}_{slide}_{idx}_{version}".encode()).hexdigest()


# =========================
# ✅ CHUNKING
# =========================
def chunk_documents(docs, chunk_size=500, chunk_overlap=100):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )

    final_docs = []

    for doc in docs:
        text = doc.page_content.strip()
        meta = doc.metadata.copy()

        if not text:
            continue

        prefix = f"Slide Title: {meta.get('slide_title','')}\n"

        if len(text) <= chunk_size:
            meta["chunk_id"] = 0
            meta["chunk_uid"] = chunk_id(meta["file_name"], meta["slide_number"], 0, meta["version"])
            final_docs.append(Document(page_content=prefix + text, metadata=meta))
            continue

        chunks = splitter.split_text(text)

        for idx, chunk in enumerate(chunks):
            if not chunk.strip():
                continue

            meta_copy = meta.copy()
            meta_copy["chunk_id"] = idx
            meta_copy["chunk_uid"] = chunk_id(meta["file_name"], meta["slide_number"], idx, meta["version"])

            final_docs.append(Document(page_content=prefix + chunk, metadata=meta_copy))

    return final_docs


# =========================
# ✅ EMBEDDING
# =========================
def get_embedding():
    return AzureOpenAIEmbeddings(
        azure_deployment=EMBEDDING_DEPLOYMENT,
        openai_api_version=API_VERSION,
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY
    )


# =========================
# ✅ VECTOR STORE
# =========================
def get_vector_store(embedding):
    return Chroma(
        collection_name="ppt_rag",
        embedding_function=embedding,
        persist_directory=PERSIST_DIR
    )


# =========================
# ✅ UPSERT
# =========================
def upsert(vs, docs):
    vs.add_texts(
        texts=[d.page_content for d in docs],
        metadatas=[d.metadata for d in docs],
        ids=[d.metadata["chunk_uid"] for d in docs]
    )


# =========================
# ✅ QUERY CLEANING
# =========================
def clean_query(query):
    return query.lower().replace("what are the comments in", "").strip()


# =========================
# ✅ RETRIEVE (FIXED LOGIC)
# =========================
def retrieve(vs, query, k=5):

    query_clean = clean_query(query)
    query_norm = normalize(query_clean)

    all_data = vs.get()

    exact = []
    partial = []

    for content, meta in zip(all_data["documents"], all_data["metadatas"]):
        title = meta.get("slide_title") or ""
        title_norm = normalize(title)

        if query_norm == title_norm:
            exact.append(Document(page_content=content, metadata=meta))

        elif query_norm in title_norm:
            partial.append(Document(page_content=content, metadata=meta))

    # ✅ priority
    if exact:
        return exact[:k]

    if partial:
        return partial[:k]

    # fallback
    return vs.similarity_search(query, k=k)


# =========================
# ✅ FILTER SAME SLIDE
# =========================
def filter_same_slide(docs):
    if not docs:
        return docs

    slide = docs[0].metadata["slide_number"]
    return [d for d in docs if d.metadata["slide_number"] == slide]


# =========================
# ✅ CONTEXT
# =========================
def build_context(docs):
    return "\n\n".join([
        f"File: {d.metadata['file_name']} | Slide: {d.metadata['slide_number']}\n{d.page_content}"
        for d in docs
    ])


# =========================
# ✅ INTENT
# =========================
def classify_query(query):
    q = query.lower()

    if "comment" in q:
        return "comments"
    if any(w in q for w in ["trend", "growth", "analysis"]):
        return "chart"
    return "general"


# =========================
# ✅ PROMPT
# =========================
def build_prompt(query, context, intent):

    if intent == "comments":
        return f"""
Extract bullet point comments.

Context:
{context}

Question:
{query}
"""

    elif intent == "chart":
        return f"""
Analyze trends.

Context:
{context}

Question:
{query}
"""

    else:
        return f"""
Answer based on context.

Context:
{context}

Question:
{query}
"""


# =========================
# ✅ LLM
# =========================
def get_llm():
    return AzureChatOpenAI(
        azure_deployment=CHAT_DEPLOYMENT,
        openai_api_version=API_VERSION,
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        temperature=0
    )


# =========================
# ✅ QA
# =========================
def answer(query, vs, llm):

    intent = classify_query(query)

    docs = retrieve(vs, query, k=8)
    docs = filter_same_slide(docs)

    if not docs:
        return "No data found."

    context = build_context(docs)

    prompt = build_prompt(query, context, intent)

    response = llm.invoke(prompt).content

    refs = "\n".join({
        f"(File: {d.metadata['file_name']}, Slide: {d.metadata['slide_number']})"
        for d in docs
    })

    return response + "\n\nSources:\n" + refs


# =========================
# ✅ MAIN
# =========================
if __name__ == "__main__":

    ppt_file = "input.pptx"

    raw = extract_ppt(ppt_file)
    docs = to_docs(raw)
    chunks = chunk_documents(docs)

    emb = get_embedding()
    vs = get_vector_store(emb)

    upsert(vs, chunks)

    print("✅ Pipeline Ready")

    llm = get_llm()

    query = "What are the comments in G&A Evaluation - MTD vs FC7+5"

    result = answer(query, vs, llm)

    print("\n=== ANSWER ===\n")
    print(result)
