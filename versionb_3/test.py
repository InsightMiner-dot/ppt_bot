import os
import hashlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter

# =========================
# ✅ CONFIG
# =========================
IMAGE_DIR = "images"


# =========================
# ✅ VERSION (FILE HASH)
# =========================
def get_file_hash(file_path):
    with open(file_path, "rb") as f:
        return hashlib.md5(f.read()).hexdigest()


# =========================
# ✅ EXTRACTION (TEXT + TABLES + NOTES + IMAGES + CHARTS)
# =========================
def extract_ppt(file_path):
    prs = Presentation(file_path)
    file_name = os.path.basename(file_path)
    version = get_file_hash(file_path)

    os.makedirs(IMAGE_DIR, exist_ok=True)

    documents = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        images = []
        charts = []
        slide_title = None

        for shape in slide.shapes:

            # -------------------------
            # 1. TEXT FRAMES
            # -------------------------
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_text.append(text)
                    # Use the first text block as the title fallback if not explicitly set
                    if not slide_title:
                        slide_title = text.split('\n')[0]

            # -------------------------
            # 2. GROUPED SHAPES
            # -------------------------
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.has_text_frame:
                        text = sub_shape.text.strip()
                        if text:
                            slide_text.append(text)

            # -------------------------
            # 3. TABLES
            # -------------------------
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    slide_text.append(" | ".join(row_data))

            # -------------------------
            # 4. IMAGES
            # -------------------------
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                ext = image.ext

                image_name = f"{file_name}_slide_{i+1}_{len(images)}{ext}"
                image_path = os.path.join(IMAGE_DIR, image_name)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                images.append({
                    "image_path": image_path,
                    "width": int(shape.width),
                    "height": int(shape.height)
                })

            # -------------------------
            # 5. CHARTS
            # -------------------------
            if shape.has_chart:
                chart = shape.chart

                chart_data = {
                    "chart_title": None,
                    "series": []
                }

                if chart.has_title:
                    chart_data["chart_title"] = chart.chart_title.text_frame.text

                for series in chart.series:
                    series_data = {"name": series.name}
                    
                    # 🔥 Map categories to values for LLM readability
                    try:
                        categories = [c.label for c in chart.plots[0].categories]
                        paired_data = dict(zip(categories, series.values))
                        series_data["values"] = paired_data
                    except Exception:
                        series_data["values"] = list(series.values)

                    chart_data["series"].append(series_data)

                charts.append(chart_data)

                # Convert chart → text (highly important for RAG)
                summary = []
                if chart_data["chart_title"]:
                    summary.append(f"Chart: {chart_data['chart_title']}")

                for s in chart_data["series"]:
                    summary.append(f"{s['name']} data: {s['values']}")

                if summary:
                    slide_text.append("\n".join(summary))

        # -------------------------
        # 6. SPEAKER NOTES
        # -------------------------
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_text.append(f"Speaker Notes: {notes_text}")

        # Assemble the raw document for the slide
        documents.append({
            # Join all elements with a double newline for clean separation
            "content": "\n\n".join(slide_text),
            "metadata": {
                "file_name": file_name,
                "version": version,
                "slide_number": i + 1,
                "slide_title": slide_title or f"Slide {i + 1}",
                "image_count": len(images),
                "images": images,
                "chart_count": len(charts),
                "charts": charts,
                "source": f"{file_name}_slide_{i+1}"
            }
        })

    return documents


# =========================
# ✅ CONVERT TO LANGCHAIN DOCS
# =========================
def to_langchain_docs(docs):
    return [
        Document(
            page_content=d["content"],
            metadata=d["metadata"]
        )
        for d in docs
    ]


# =========================
# ✅ STABLE CHUNK ID
# =========================
def generate_chunk_id(file_name, slide_number, chunk_id, version):
    base = f"{file_name}_{slide_number}_{chunk_id}_{version}"
    return hashlib.md5(base.encode()).hexdigest()


# =========================
# ✅ SLIDE-AWARE CHUNKING
# =========================
def chunk_documents_slide_level(documents, max_chunk_size=2000, chunk_overlap=200):
    """
    1 Slide = 1 Chunk logic. 
    Only uses recursive splitting as a fallback if the slide contains an massive amount of text.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=max_chunk_size,
        chunk_overlap=chunk_overlap
    )

    chunked_docs = []

    for doc in documents:
        text = doc.page_content.strip()
        metadata = doc.metadata.copy()

        if not text:
            continue

        # 🔥 Global + Local Context Header Injection
        file_name = metadata.get("file_name", "Unknown File")
        slide_title = metadata.get("slide_title", "Untitled Slide")
        slide_num = metadata.get("slide_number", "?")
        
        context_header = f"[Presentation: {file_name} | Slide {slide_num}: {slide_title}]\n\n"
        full_text = context_header + text

        # -------------------------
        # IDEAL CASE: Whole slide stays as one chunk
        # -------------------------
        if len(full_text) <= max_chunk_size:
            metadata["chunk_id"] = 0
            metadata["chunk_uid"] = generate_chunk_id(
                file_name, slide_num, 0, metadata["version"]
            )
            
            chunked_docs.append(
                Document(page_content=full_text, metadata=metadata)
            )
            continue

        # -------------------------
        # FALLBACK CASE: Slide is massive (split it)
        # -------------------------
        # Split only the body text, so we can re-attach the context header to every piece
        chunks = splitter.split_text(text)
        
        for idx, chunk in enumerate(chunks):
            chunk = chunk.strip()
            
            metadata_copy = metadata.copy()
            metadata_copy["chunk_id"] = idx
            metadata_copy["chunk_uid"] = generate_chunk_id(
                file_name, slide_num, idx, metadata["version"]
            )

            # Re-inject the context header into EVERY sub-chunk
            chunked_docs.append(
                Document(
                    page_content=context_header + chunk, 
                    metadata=metadata_copy
                )
            )

    return chunked_docs


# =========================
# ✅ MAIN (NOTEBOOK / SCRIPT)
# =========================
if __name__ == "__main__":
    ppt_file = "input.pptx"

    # Ensure the file exists before running
    if not os.path.exists(ppt_file):
        print(f"Error: {ppt_file} not found. Please provide a valid file.")
    else:
        # Step 1: Extraction
        raw_docs = extract_ppt(ppt_file)

        # Step 2: Convert
        lc_docs = to_langchain_docs(raw_docs)

        # Step 3: Slide-Level Chunking
        chunked_docs = chunk_documents_slide_level(lc_docs)

        # Debug output
        print(f"\n✅ Pipeline Complete! Total chunks generated: {len(chunked_docs)}\n")

        for d in chunked_docs[:3]:
            print(f"--- Chunk UID: {d.metadata['chunk_uid']} ---")
            print(f"Content Preview:\n{d.page_content[:250]}...")
            print(f"Metadata Keys: {list(d.metadata.keys())}\n")
