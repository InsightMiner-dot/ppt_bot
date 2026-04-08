import os
import hashlib
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ✅ Modern LangChain Imports
from langchain_core.documents import Document
from langchain_text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import AzureOpenAIEmbeddings
from langchain_chroma import Chroma

# Load environment variables
load_dotenv()

# =========================
# ✅ CONFIGURATION
# =========================
IMAGE_DIR = "images"
CHROMA_PERSIST_DIR = "./chroma_db_storage"
COLLECTION_NAME = "ppt_production_kb"
AZURE_DEPLOYMENT_NAME = "text-embedding-3-small" # Change to your actual deployment name

# =========================
# ✅ UTILITIES
# =========================
def get_file_hash(file_path):
    """Generates an MD5 hash of the file for version control."""
    with open(file_path, "rb") as f:
        return hashlib.md5(f.read()).hexdigest()

def generate_chunk_id(file_name, slide_number, chunk_id, version):
    """Generates a stable, deterministic ID for vector DB deduplication."""
    base = f"{file_name}_{slide_number}_{chunk_id}_{version}"
    return hashlib.md5(base.encode()).hexdigest()

# =========================
# ✅ 1. EXTRACTION
# =========================
def extract_ppt(file_path):
    """Extracts text, tables, notes, grouped shapes, charts, and images from a PPTX."""
    prs = Presentation(file_path)
    file_name = os.path.basename(file_path)
    version = get_file_hash(file_path)

    os.makedirs(IMAGE_DIR, exist_ok=True)
    documents = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        images = []
        charts = []
        slide_title = None

        for shape in slide.shapes:
            # 1. TEXT FRAMES
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_text.append(text)
                    if not slide_title:
                        slide_title = text.split('\n')[0]

            # 2. GROUPED SHAPES
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.has_text_frame:
                        text = sub_shape.text.strip()
                        if text:
                            slide_text.append(text)

            # 3. TABLES
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    slide_text.append(" | ".join(row_data))

            # 4. IMAGES
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                image_name = f"{file_name}_slide_{i+1}_{len(images)}{ext}"
                image_path = os.path.join(IMAGE_DIR, image_name)

                with open(image_path, "wb") as f:
                    f.write(image.blob)

                images.append({
                    "image_path": image_path,
                    "width": int(shape.width),
                    "height": int(shape.height)
                })

            # 5. CHARTS
            if shape.has_chart:
                chart = shape.chart
                chart_data = {"chart_title": None, "series": []}

                if chart.has_title:
                    chart_data["chart_title"] = chart.chart_title.text_frame.text

                for series in chart.series:
                    series_data = {"name": series.name}
                    try:
                        categories = [c.label for c in chart.plots[0].categories]
                        series_data["values"] = dict(zip(categories, series.values))
                    except Exception:
                        series_data["values"] = list(series.values)
                    chart_data["series"].append(series_data)

                charts.append(chart_data)

                # Convert chart to text for vector DB
                summary = []
                if chart_data["chart_title"]:
                    summary.append(f"Chart: {chart_data['chart_title']}")
                for s in chart_data["series"]:
                    summary.append(f"{s['name']} data: {s['values']}")
                if summary:
                    slide_text.append("\n".join(summary))

        # 6. SPEAKER NOTES
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_text.append(f"Speaker Notes: {notes_text}")

        # Assemble Raw Document
        documents.append({
            "content": "\n\n".join(slide_text),
            "metadata": {
                "file_name": file_name,
                "version": version,
                "slide_number": i + 1,
                "slide_title": slide_title or f"Slide {i + 1}",
                "image_count": len(images),
                "images": str(images), # Stringified for Chroma compatibility
                "chart_count": len(charts),
                "source": f"{file_name}_slide_{i+1}"
            }
        })

    return documents

# =========================
# ✅ 2. CHUNKING
# =========================
def chunk_documents_slide_level(raw_docs, max_chunk_size=2000, chunk_overlap=200):
    """Implements 1 Slide = 1 Chunk with global context injection."""
    
    # Convert dicts to LangChain Document objects first
    lc_docs = [Document(page_content=d["content"], metadata=d["metadata"]) for d in raw_docs]
    
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=max_chunk_size,
        chunk_overlap=chunk_overlap
    )

    chunked_docs = []

    for doc in lc_docs:
        text = doc.page_content.strip()
        metadata = doc.metadata.copy()

        if not text:
            continue

        file_name = metadata.get("file_name", "Unknown File")
        slide_title = metadata.get("slide_title", "Untitled Slide")
        slide_num = metadata.get("slide_number", "?")
        
        # Inject context header so the LLM knows WHERE this chunk came from
        context_header = f"[Presentation: {file_name} | Slide {slide_num}: {slide_title}]\n\n"
        full_text = context_header + text

        # IDEAL CASE: Whole slide is one chunk
        if len(full_text) <= max_chunk_size:
            metadata["chunk_id"] = 0
            metadata["chunk_uid"] = generate_chunk_id(file_name, slide_num, 0, metadata["version"])
            chunked_docs.append(Document(page_content=full_text, metadata=metadata))
            continue

        # FALLBACK CASE: Slide is massive (split body, re-inject header)
        chunks = splitter.split_text(text)
        for idx, chunk in enumerate(chunks):
            metadata_copy = metadata.copy()
            metadata_copy["chunk_id"] = idx
            metadata_copy["chunk_uid"] = generate_chunk_id(file_name, slide_num, idx, metadata["version"])
            chunked_docs.append(Document(page_content=context_header + chunk.strip(), metadata=metadata_copy))

    return chunked_docs

# =========================
# ✅ 3. INGESTION
# =========================
def ingest_into_chroma(documents):
    """Embeds and stores chunks in ChromaDB using Azure OpenAI."""
    if not documents:
        print("⚠️ No documents to ingest.")
        return None

    print(f"\n🚀 Initializing Azure OpenAI embeddings...")
    embedding_model = AzureOpenAIEmbeddings(
        azure_deployment=AZURE_DEPLOYMENT_NAME
    )

    print(f"📦 Storing {len(documents)} chunks in ChromaDB at {CHROMA_PERSIST_DIR}...")
    vector_store = Chroma.from_documents(
        documents=documents,
        embedding=embedding_model,
        collection_name=COLLECTION_NAME,
        persist_directory=CHROMA_PERSIST_DIR,
    )
    
    print(f"✅ Successfully persisted vector database!")
    return vector_store

# =========================
# ✅ MAIN EXECUTION
# =========================
if __name__ == "__main__":
    ppt_file = "input.pptx" # REPLACE WITH YOUR FILE PATH

    if not os.path.exists(ppt_file):
        print(f"❌ Error: {ppt_file} not found. Please provide a valid file.")
    else:
        # Step 1: Extract
        print(f"📄 Extracting content from {ppt_file}...")
        raw_documents = extract_ppt(ppt_file)

        # Step 2: Chunk
        print(f"✂️  Applying slide-level chunking...")
        chunked_documents = chunk_documents_slide_level(raw_documents)
        print(f"   -> Generated {len(chunked_documents)} total chunks.")

        # Step 3: Embed & Ingest
        try:
            vector_db = ingest_into_chroma(chunked_documents)
            
            # Step 4: Quick Test Query
            if vector_db:
                print("\n🔍 Running a quick test search...")
                # Change this query to something that actually exists in your PPTX
                query = "What are the key takeaways?" 
                
                results = vector_db.similarity_search_with_score(query, k=2)
                
                for doc, score in results:
                    print(f"\n--- Distance Score: {score:.4f} ---")
                    print(f"Metadata: Slide {doc.metadata.get('slide_number')} | {doc.metadata.get('slide_title')}")
                    print(f"Content Preview: {doc.page_content[:150]}...")

        except Exception as e:
            print(f"\n❌ Error during ingestion/querying: {e}")
            print("Check your .env file and Azure deployment names.")
