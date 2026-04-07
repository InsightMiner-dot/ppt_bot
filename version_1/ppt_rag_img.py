import os
import sys
import logging
import glob
import base64
from pptx import Presentation
from scripts import llm  # Make sure this imports your llm.py file

# --- 1. Logging Setup ---
log_dir = "logs"
os.makedirs(log_dir, exist_ok=True)
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler(os.path.join(log_dir, "process.log"), encoding="utf-8"),
        logging.StreamHandler(sys.stdout)
    ]
)

logging.info("Multimodal Multi-PPT Script started.")

# --- 2. Find Presentations ---
target_folder = "presentations"
os.makedirs(target_folder, exist_ok=True) 
ppt_files = glob.glob(os.path.join(target_folder, "*.pptx"))

if not ppt_files:
    logging.warning(f"No .pptx files found in '{target_folder}'. Please add some and run again.")
    sys.exit(0)

logging.info(f"Found {len(ppt_files)} presentations to process.")

# --- 3. Single-Pass Extraction (Text, Tables, Images) ---
ppt_data = {}

for file_path in ppt_files:
    try:
        logging.info(f"Processing: {file_path}")
        actual_filename = os.path.basename(file_path)
        prs = Presentation(file_path)
        
        # Loop through exactly 1, 2, 3...
        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_content_elements = []
            
            for shape in slide.shapes:
                # A. Handle Standard Text
                if shape.has_text_frame:
                    if shape.text.strip():
                        slide_content_elements.append(shape.text.strip())
                
                # B. Handle Native Tables (Cheaper/Faster than sending a table image to Vision)
                elif shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        # Clean up cell text and separate with a pipe '|'
                        row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                        table_rows.append(" | ".join(row_data))
                    
                    if table_rows:
                        slide_content_elements.append("### Table Data ###\n" + "\n".join(table_rows))
                
                # C. Handle Images and Charts (Vision LLM)
                elif hasattr(shape, "image"):
                    logging.info(f"   -> Sending image from Slide {slide_index} to Vision LLM...")
                    base64_img = base64.b64encode(shape.image.blob).decode('utf-8')
                    image_metadata = {
                        "filename": actual_filename,
                        "page_number": slide_index
                    }
                    try:
                        desc = llm.describe_image(base64_image=base64_img, metadata=image_metadata)
                        slide_content_elements.append(f"### Image/Chart Description ###\n{desc}")
                    except Exception as img_err:
                        logging.error(f"   -> Vision LLM failed for image on slide {slide_index}: {img_err}")
            
            # If the slide actually had content, save it
            if slide_content_elements:
                slide_text = "\n\n".join(slide_content_elements)
                
                page = str(slide_index) 
                source = file_path
                filename = actual_filename
                last_modified = "Unknown Date" 
                
                unique_key = f"{filename}_Page_{page}"
                meta_string = f"[Reference Metadata -> Source: {source}, Filename: {filename}, Last Modified: {last_modified}, Page Number: {page}]\n"
                
                ppt_data[unique_key] = f"{meta_string}\n{slide_text}"
                
    except Exception as e:
        logging.error(f"Failed to process {file_path}. Skipping. Error: {e}")

logging.info(f"Extraction complete. Total slides successfully mapped: {len(ppt_data)}")

# --- 4. Build Final Context String ---
context = ""
for unique_key, content in ppt_data.items():
    context += f"## {unique_key}:\n\n {content.strip()}\n\n\n"

# --- 5. Query the LLM ---
question = "Summarize the key data, including any trends from the charts or tables, across all the provided presentations."

# Strict instruction to enforce the bulleted file references at the end
instruction = """

IMPORTANT INSTRUCTION FOR REFERENCES:
You are summarizing data across multiple different files. At the very end of your response, you MUST include a comprehensive "References" section. 
You must list EVERY single file and page number that contributed to your answer. Format it as a bulleted list exactly like this:
- Filename: [Insert Filename], Page Number: [Insert Page Number]
- Filename: [Insert Filename], Page Number: [Insert Page Number]
"""

try:
    logging.info("Sending combined Text, Table, and Vision context to Azure OpenAI...")
    response = llm.ask_llm(question=question + instruction, context=context)
    
    logging.info("Response received successfully.")
    print("\n" + "="*50)
    print("LLM RESPONSE:")
    print("="*50)
    print(response)
    
    with open(r"response.md", "w", encoding="utf-8") as f:
        f.write(response)
    logging.info("Output successfully saved to response.md")

except Exception as e:
    logging.error(f"Error querying the final LLM: {e}")

logging.info("Script finished.")