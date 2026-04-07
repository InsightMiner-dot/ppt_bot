import base64
import glob
import hashlib
import io
import logging
import os
import shutil
import sys
from pathlib import Path

from langchain_chroma import Chroma
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation

CURRENT_DIR = Path(__file__).resolve().parent
if str(CURRENT_DIR) not in sys.path:
    sys.path.insert(0, str(CURRENT_DIR))

from scripts import llm

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")


def _clean_text(value: str) -> str:
    return " ".join(value.split())


def _slide_header(filename: str, slide_index: int) -> str:
    return f"Filename: {filename}\nSlide Number: {slide_index}"


def _pick_slide_title(text_blocks: list[str]) -> str:
    for block in text_blocks:
        for line in block.splitlines():
            cleaned_line = _clean_text(line)
            if cleaned_line:
                return cleaned_line
    return ""


def _extract_chart_text(shape, filename: str, slide_index: int, chart_index: int) -> str | None:
    try:
        chart = shape.chart
        chart_parts = [
            _slide_header(filename, slide_index),
            "Content Type: Chart",
            f"Chart Index: {chart_index}",
        ]

        chart_title = None
        if chart.has_title and chart.chart_title and chart.chart_title.text_frame:
            chart_title = _clean_text(chart.chart_title.text_frame.text)
        if chart_title:
            chart_parts.append(f"Chart Title: {chart_title}")

        category_labels = []
        try:
            categories = chart.plots[0].categories
            for category in categories:
                label = _clean_text(str(category.label))
                if label:
                    category_labels.append(label)
        except Exception:
            category_labels = []

        if category_labels:
            chart_parts.append("Categories: " + " | ".join(category_labels))

        series_lines = []
        for series in chart.series:
            series_name = _clean_text(getattr(series, "name", "") or "Unnamed Series")
            points = []
            values = list(series.values)
            for idx, value in enumerate(values):
                label = category_labels[idx] if idx < len(category_labels) else f"Point {idx + 1}"
                points.append(f"{label}: {value}")
            series_lines.append(f"{series_name} -> " + "; ".join(points))

        if series_lines:
            chart_parts.append("Series Data:\n" + "\n".join(series_lines))

        if len(chart_parts) <= 3:
            return None

        return "\n".join(chart_parts)
    except Exception as exc:
        logging.warning(f"Chart extraction failed on slide {slide_index}: {exc}")
        return None


def _extract_ocr_text(image_blob: bytes, slide_index: int, image_index: int) -> str:
    try:
        from PIL import Image
        import pytesseract
    except ImportError:
        logging.info(
            "OCR dependencies not available. Install 'pytesseract' and 'Pillow' to enable OCR."
        )
        return ""

    try:
        image = Image.open(io.BytesIO(image_blob))
        extracted_text = pytesseract.image_to_string(image)
        cleaned_text = _clean_text(extracted_text)
        if cleaned_text:
            logging.info(
                f"   -> OCR extracted text for image {image_index} on Slide {slide_index}."
            )
        return cleaned_text
    except Exception as exc:
        logging.warning(
            f"OCR failed for image {image_index} on slide {slide_index}: {exc}"
        )
        return ""


def _make_document(
    file_path: str,
    filename: str,
    slide_index: int,
    content_type: str,
    content: str,
    slide_title: str = "",
    element_index: int | None = None,
) -> Document:
    metadata = {
        "source": file_path,
        "filename": filename,
        "page": str(slide_index),
        "content_type": content_type,
        "slide_title": slide_title,
    }
    if element_index is not None:
        metadata["element_index"] = element_index

    return Document(page_content=content, metadata=metadata)


def build_vector_database():
    target_folder = "presentations"
    os.makedirs(target_folder, exist_ok=True)
    ppt_files = glob.glob(os.path.join(target_folder, "*.pptx"))

    if not ppt_files:
        logging.warning("No .pptx files found. Please add some and run again.")
        return

    raw_documents: list[Document] = []

    for file_path in ppt_files:
        actual_filename = os.path.basename(file_path)
        logging.info(f"Extracting: {actual_filename}")

        try:
            prs = Presentation(file_path)
            for slide_index, slide in enumerate(prs.slides, start=1):
                text_blocks = []
                table_blocks = []
                chart_blocks = []
                image_blocks = []
                seen_image_hashes = set()
                table_counter = 0
                chart_counter = 0
                image_counter = 0

                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        cleaned_text = _clean_text(shape.text)
                        if cleaned_text:
                            text_blocks.append(cleaned_text)
                        continue

                    slide_title = _pick_slide_title(text_blocks)

                    if shape.has_table:
                        table_rows = []
                        for row in shape.table.rows:
                            row_data = [_clean_text(cell.text) for cell in row.cells]
                            table_rows.append(" | ".join(row_data))

                        if table_rows:
                            table_counter += 1
                            table_text = (
                                f"{_slide_header(actual_filename, slide_index)}\n"
                                "Content Type: Table\n"
                                f"Table Index: {table_counter}\n"
                                "Table Data:\n"
                                + "\n".join(table_rows)
                            )
                            table_blocks.append(table_text)
                            raw_documents.append(
                                _make_document(
                                    file_path=file_path,
                                    filename=actual_filename,
                                    slide_index=slide_index,
                                    content_type="table",
                                    content=table_text,
                                    slide_title=slide_title,
                                    element_index=table_counter,
                                )
                            )
                        continue

                    if getattr(shape, "has_chart", False):
                        chart_counter += 1
                        chart_text = _extract_chart_text(
                            shape=shape,
                            filename=actual_filename,
                            slide_index=slide_index,
                            chart_index=chart_counter,
                        )
                        if chart_text:
                            chart_blocks.append(chart_text)
                            raw_documents.append(
                                _make_document(
                                    file_path=file_path,
                                    filename=actual_filename,
                                    slide_index=slide_index,
                                    content_type="chart",
                                    content=chart_text,
                                    slide_title=slide_title,
                                    element_index=chart_counter,
                                )
                            )
                        continue

                    if hasattr(shape, "image"):
                        image_blob = shape.image.blob
                        image_hash = hashlib.sha256(image_blob).hexdigest()
                        if image_hash in seen_image_hashes:
                            continue

                        seen_image_hashes.add(image_hash)
                        image_counter += 1
                        logging.info(
                            f"   -> Processing image {image_counter} on Slide {slide_index}..."
                        )
                        base64_img = base64.b64encode(image_blob).decode("utf-8")
                        ocr_text = _extract_ocr_text(
                            image_blob=image_blob,
                            slide_index=slide_index,
                            image_index=image_counter,
                        )
                        image_metadata = {
                            "filename": actual_filename,
                            "page_number": slide_index,
                            "slide_text": "\n".join(text_blocks),
                            "chart_text": "\n".join(chart_blocks),
                            "ocr_text": ocr_text,
                            "slide_title": slide_title,
                        }
                        try:
                            desc = llm.describe_image(
                                base64_image=base64_img,
                                metadata=image_metadata,
                            )
                            cleaned_desc = _clean_text(desc)
                            if cleaned_desc:
                                image_sections = []
                                if ocr_text:
                                    image_sections.append(f"OCR Text:\n{ocr_text}")
                                image_sections.append(f"Image Description:\n{cleaned_desc}")
                                image_text = (
                                    f"{_slide_header(actual_filename, slide_index)}\n"
                                    "Content Type: Image\n"
                                    f"Image Index: {image_counter}\n"
                                    + "\n\n".join(image_sections)
                                )
                                image_blocks.append(image_text)
                                raw_documents.append(
                                    _make_document(
                                        file_path=file_path,
                                        filename=actual_filename,
                                        slide_index=slide_index,
                                        content_type="image",
                                        content=image_text,
                                        slide_title=slide_title,
                                        element_index=image_counter,
                                    )
                                )
                        except Exception as exc:
                            logging.error(f"Vision failed on slide {slide_index}: {exc}")

                slide_sections = []
                if text_blocks:
                    slide_sections.append("Slide Text:\n" + "\n".join(text_blocks))
                if table_blocks:
                    slide_sections.append("Tables:\n" + "\n\n".join(table_blocks))
                if chart_blocks:
                    slide_sections.append("Charts:\n" + "\n\n".join(chart_blocks))
                if image_blocks:
                    slide_sections.append("Images:\n" + "\n\n".join(image_blocks))

                if slide_sections:
                    slide_title = _pick_slide_title(text_blocks)
                    slide_text = (
                        f"{_slide_header(actual_filename, slide_index)}\n"
                        "Content Type: Slide Summary\n\n"
                        + "\n\n".join(slide_sections)
                    )
                    raw_documents.append(
                        _make_document(
                            file_path=file_path,
                            filename=actual_filename,
                            slide_index=slide_index,
                            content_type="slide_summary",
                            content=slide_text,
                            slide_title=slide_title,
                        )
                    )

        except Exception as exc:
            logging.error(f"Failed to process {file_path}: {exc}")

    logging.info(f"Extracted {len(raw_documents)} total retrievable items. Starting chunking...")

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1500,
        chunk_overlap=150,
    )
    chunked_docs = text_splitter.split_documents(raw_documents)
    logging.info(f"Created {len(chunked_docs)} chunks.")

    db_directory = "./chroma_db"
    logging.info("Saving to Chroma DB...")

    if os.path.isdir(db_directory):
        shutil.rmtree(db_directory)

    Chroma.from_documents(
        documents=chunked_docs,
        embedding=llm.embeddings,
        persist_directory=db_directory,
    )

    logging.info("Database built successfully!")


if __name__ == "__main__":
    build_vector_database()
